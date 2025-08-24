# app.py for Hugging Face Spaces - RAG File Vault

# --- 1. Standard Imports ---
import os
import re
import time
import pandas as pd
from tqdm import tqdm
from bs4 import BeautifulSoup
import fitz  # PyMuPDF
import pdfplumber
import docx
import pptx
import numpy as np
import logging
import uuid # For unique point IDs
import zipfile # MODIFICATION: Added to validate docx files

# --- 2. Cloud Service and AI Library Imports ---
import google.generativeai as genai
from supabase import create_client, Client
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct
from qdrant_client.http import models as rest
from langchain.text_splitter import RecursiveCharacterTextSplitter
from google.generativeai import GenerativeModel # For the chat model

# --- 3. Gradio Import ---
import gradio as gr

# Suppress verbose warnings from pdfminer used by pdfplumber
logging.getLogger("pdfminer").setLevel(logging.ERROR)

# --- 4. Configuration and Constants ---
BUCKET_NAME = "rag-file"
COLLECTION_NAME = "rag_embeddings"
BATCH_SIZE = 128 # Qdrant upsert batch size
EMBED_MODEL = "models/embedding-001"
CHAT_MODEL_NAME = "gemini-1.5-flash"
SUPABASE_TOTAL_BYTES = 1073741824  # 1 GB
QDRANT_MAX_VECTORS = 100000
DEFAULT_EMBEDDING_BATCH_SIZE = 50

# --- 5. Global Client Placeholders ---
supabase: Client = None
qdrant = None

# --- 6. Core Functions (Extract, Embed, Store, Delete, Query, Setup) ---

def extract_text_from_file(file_path):
    """Extracts text from various file types with a robust PDF extraction strategy."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            text = ""
            try:
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text()
                if text.strip():
                    print("ℹ️ Extracted text using PyMuPDF.")
                    return text.strip()
            except Exception as fitz_e:
                print(f"⚠️ PyMuPDF failed: {fitz_e}. Falling back to pdfplumber.")
            text = ""
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        text += page_text
                print("ℹ️ Extracted text using pdfplumber.")
                return text.strip()
            except Exception as pdf_e:
                return f"[PDF Extraction Error: Both libraries failed. Last error: {pdf_e}]"
        elif ext == ".docx":
            # MODIFICATION START: Validate that the .docx is a proper zip archive before processing.
            if not zipfile.is_zipfile(file_path):
                return "[Extraction Error: The file is not a valid .docx (Zip archive). It might be a legacy .doc file renamed to .docx, or the file is corrupt.]"
            # If valid, proceed with the original, direct method.
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs]).strip()
            # MODIFICATION END
        elif ext == ".pptx":
            prs = pptx.Presentation(file_path)
            text_runs = [run.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
            return "\n".join(text_runs).strip()
        elif ext in [".csv", ".xlsx"]:
            df = pd.read_csv(file_path, on_bad_lines='skip') if ext == ".csv" else pd.read_excel(file_path)
            return df.to_string(index=False).strip()
        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()
        elif ext == ".html":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return BeautifulSoup(f, "html.parser").get_text(separator="\n").strip()
        else:
            return f"[Unsupported file type: {ext}]"
    except Exception as e:
        return f"[Error extracting text from {ext} file: {str(e)}]"

def embed_text(text, chunk_size=1000, batch_size=DEFAULT_EMBEDDING_BATCH_SIZE):
    """Splits text and generates embeddings using batching."""
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=100)
    chunks = splitter.split_text(text)
    chunks = [c for c in chunks if c.strip()]
    if not chunks: return np.array([]), []
    embeddings = []
    for i in tqdm(range(0, len(chunks), batch_size), desc="Embedding Chunks"):
        batch_chunks = chunks[i:i + batch_size]
        try:
            result = genai.embed_content(model=EMBED_MODEL, content=batch_chunks)
            embeddings.extend(result['embedding'])
        except Exception as e:
            print(f"⚠️ Error embedding batch at index {i}: {e}")
            embeddings.extend([[0.0] * 768] * len(batch_chunks))
    return np.array(embeddings), chunks

def sanitise_key(name: str) -> str:
    """Sanitizes filename for safe storage."""
    return re.sub(r'[^0-9A-Za-z._-]+', '-', name)

def store_file_and_embeddings(supabase_client, qdrant_client, filename, file_bytes, embeddings, chunks):
    """Uploads file to Supabase and embeddings to Qdrant, handling batching."""
    storage_path, file_url, supabase_error, qdrant_error = filename, "", None, None
    try:
        supabase_url = os.environ.get('SUPABASE_URL', '')
        supabase_client.storage.from_(BUCKET_NAME).upload(path=storage_path, file=file_bytes, file_options={"upsert": "true"})
        file_url = f"{supabase_url}/storage/v1/object/public/{BUCKET_NAME}/{storage_path}"
        print(f"✅ File stored at: {file_url}")
    except Exception as e:
        supabase_error = f"Supabase upload failed: {e}"
        print(f"⚠️ {supabase_error} for {filename}")
    if embeddings.size > 0:
        points = [PointStruct(id=str(uuid.uuid4()), vector=v.tolist(), payload={"file_name": filename, "chunk_id": i, "text": c, "file_url": file_url}) for i, (v, c) in enumerate(zip(embeddings, chunks))]
        for i in tqdm(range(0, len(points), BATCH_SIZE), desc="Upserting to Qdrant"):
            try:
                qdrant_client.upsert(collection_name=COLLECTION_NAME, points=points[i:i + BATCH_SIZE], wait=True)
            except Exception as e:
                qdrant_error = f"Qdrant upsert failed on batch {i//BATCH_SIZE}: {e}"
                print(f"⚠️ {qdrant_error}")
                break
    else:
        print(f"ℹ️ No embeddings to store for {filename}")
    if supabase_error and qdrant_error: return False, f"❌ Failed: {supabase_error}; {qdrant_error}"
    if supabase_error: return False, f"⚠️ Partial Success (Vectors may be stored): {supabase_error}"
    if qdrant_error: return False, f"⚠️ Partial Success (File stored): {qdrant_error}"
    return True, "✅ Successfully processed and stored."

def ensure_keyword_index(qdrant_client, collection_name, field_name="file_name"):
    """Creates a keyword index on a Qdrant payload field if it doesn't exist."""
    try:
        qdrant_client.create_payload_index(collection_name=collection_name, field_name=field_name, field_schema="keyword")
        print(f"✅ Ensured keyword index exists on '{field_name}'.")
    except Exception as e:
        print(f"ℹ️ Note on keyword index creation: {e} (Often safe if index already exists).")

def delete_file(supabase_client, qdrant_client, filename, delete_vectors=True, delete_file=True):
    """Deletes a file from Supabase and/or its embeddings from Qdrant."""
    log_messages = []
    if delete_vectors:
        try:
            ensure_keyword_index(qdrant_client, COLLECTION_NAME, "file_name")
            qdrant_client.delete(collection_name=COLLECTION_NAME, points_selector=rest.FilterSelector(filter=rest.Filter(must=[rest.FieldCondition(key="file_name", match=rest.MatchValue(value=filename))])))
            msg = f"✅ Deleted embeddings for '{filename}'."
            print(msg); log_messages.append(msg)
        except Exception as e:
            msg = f"⚠️ Qdrant delete error for '{filename}': {e}"
            print(msg); log_messages.append(msg)
    if delete_file:
        try:
            supabase_client.storage.from_(BUCKET_NAME).remove([filename])
            msg = f"✅ Deleted file '{filename}' from Supabase."
            print(msg); log_messages.append(msg)
        except Exception as e:
            msg = f"⚠️ Supabase delete error for '{filename}': {e}"
            print(msg); log_messages.append(msg)
    return "\n".join(log_messages)

def rag_query(qdrant_client, user_query_text: str, top_k: int = 10) -> tuple[str, str]:
    """Performs the RAG query process."""
    if not user_query_text.strip(): return "Please enter a query.", ""
    try:
        embedded_query = genai.embed_content(model=EMBED_MODEL, content=user_query_text)["embedding"]
        print("✅ Query embedded.")
        search_results = qdrant_client.search(collection_name=COLLECTION_NAME, query_vector=embedded_query, limit=top_k, with_payload=True)
        print(f"📦 Vectors returned: {len(search_results)}")
        context = "\n\n".join(r.payload.get("text", "") for r in search_results)
        prompt = f"Use the following context, combine it with your natural language processing (NLP) ability to answer the question, always make sure to simplify difficult words and concepts in simple everyday language, and provide suitable example where needed. At the end, state how helpful the context was.\nContext: {context}\nUser Question: {user_query_text}\nAnswer:"
        response = GenerativeModel(CHAT_MODEL_NAME).generate_content(prompt)
        answer = response.text
        print("✅ Answer generated.")
        formatted_context = "\n\n---\n\n".join(f"Source: {r.payload['file_name']}\n\nContent: {r.payload['text']}" for r in search_results)
        return answer, formatted_context
    except Exception as e:
        error_msg = f"Error during RAG query: {e}"
        print(f"❌ {error_msg}")
        return error_msg, ""

def setup_storage_and_vector_db(supabase_client, qdrant_client):
    """Creates Supabase bucket and Qdrant collection if they don't exist."""
    print("📦 Setting up cloud storage and vector database...")
    try:
        if not any(b.name == BUCKET_NAME for b in supabase_client.storage.list_buckets()):
            supabase_client.storage.create_bucket(BUCKET_NAME, options={"public": True})
            print(f"✅ Created Supabase bucket: {BUCKET_NAME}")
        else:
            print(f"✅ Supabase bucket '{BUCKET_NAME}' already exists.")
    except Exception as e:
        print(f"⚠️ Error with Supabase bucket setup: {e}"); raise
    try:
        if not qdrant_client.collection_exists(collection_name=COLLECTION_NAME):
            qdrant_client.create_collection(collection_name=COLLECTION_NAME, vectors_config=VectorParams(size=768, distance=Distance.COSINE))
            print(f"✅ Created Qdrant collection: {COLLECTION_NAME}")
        else:
            print(f"✅ Qdrant collection '{COLLECTION_NAME}' already exists.")
    except Exception as e:
        print(f"⚠️ Error with Qdrant collection setup: {e}"); raise

# --- 7. Gradio App Definition & Handlers ---
def create_gradio_app():
    """Defines and returns the Gradio application layout and logic."""
    def handle_credentials_submit(google_key, supabase_url, supabase_key, qdrant_url, qdrant_key):
        global supabase, qdrant
        os.environ.update({'GOOGLE_API_KEY': google_key, 'SUPABASE_URL': supabase_url, 'SUPABASE_KEY': supabase_key, 'QDRANT_URL': qdrant_url, 'QDRANT_API_KEY': qdrant_key})
        try:
            genai.configure(api_key=google_key)
            supabase = create_client(supabase_url, supabase_key)
            qdrant = QdrantClient(url=qdrant_url, api_key=qdrant_key)
            setup_storage_and_vector_db(supabase, qdrant)
            return "✅ Credentials validated and clients initialized successfully!", "", "", "", "", ""
        except Exception as e:
            return f"❌ Error initializing clients: {e}", google_key, supabase_url, supabase_key, qdrant_url, qdrant_key

    def handle_file_upload(files, progress=gr.Progress()):
        if not all((supabase, qdrant)): return "⚠️ Please set API keys first.", ""
        if not files: return "Please upload at least one file.", ""
        log, all_texts_preview, total_files = [], "", len(files)
        for i, file_obj in enumerate(files):
            filename = os.path.basename(file_obj.name)
            progress((i + 1) / total_files, desc=f"Processing: {filename}")
            log.append(f"📦 Processing '{filename}'...")
            extracted_text = extract_text_from_file(file_obj.name)
            if extracted_text.startswith(("[", "Unsupported")):
                log.append(f"⚠️ Skipped '{filename}': {extracted_text}"); continue
            all_texts_preview += f"--- Preview from {filename} ---\n{extracted_text[:500]}...\n\n"
            embeddings, chunks = embed_text(extracted_text)
            if not chunks:
                log.append(f"⚠️ No text to embed for '{filename}'."); continue
            with open(file_obj.name, "rb") as f: file_bytes = f.read()
            safe_key = sanitise_key(filename)
            success, status_msg = store_file_and_embeddings(supabase, qdrant, safe_key, file_bytes, embeddings, chunks)
            log.append(f"{status_msg} -> '{filename}'.")
        return "\n".join(log), all_texts_preview

    def handle_rag_query(query):
        if not qdrant: return "⚠️ Please set API keys first.", ""
        if not query: return "Please enter a query.", ""
        return rag_query(qdrant, query)

    def handle_check_storage():
        if not all((supabase, qdrant)): return "⚠️ Please set API keys first.", "⚠️ Please set API keys first."
        try:
            files = supabase.storage.from_(BUCKET_NAME).list()
            used_bytes = sum(f.get('metadata', {}).get('size', 0) for f in files)
            used_mb = used_bytes / (1024 * 1024)
            total_mb = SUPABASE_TOTAL_BYTES / (1024 * 1024)
            percent = (used_mb / total_mb) * 100 if total_mb > 0 else 0
            file_names = [f['name'] for f in files] if files else ["No files found."]
            supabase_status = f"=== Supabase Storage Status ===\nUsed: {used_mb:.2f} MB / {total_mb:.2f} MB ({percent:.2f}%)\n\n--- Stored Files ({len(file_names)}) ---\n" + "\n".join(file_names)
        except Exception as e: supabase_status = f"Error checking Supabase: {e}"
        try:
            info = qdrant.get_collection(collection_name=COLLECTION_NAME)
            count = info.points_count or 0
            percent = (count / QDRANT_MAX_VECTORS) * 100 if QDRANT_MAX_VECTORS > 0 else 0
            scroll_res, _ = qdrant.scroll(collection_name=COLLECTION_NAME, with_payload=["file_name"], limit=1000)
            files = sorted(list(set(p.payload['file_name'] for p in scroll_res))) or ["No embeddings found."]
            qdrant_status = f"=== Qdrant Collection Status ===\nIndexed Vectors: {count} / {QDRANT_MAX_VECTORS} ({percent:.2f}%)\n\n--- Indexed Files (Sample) ---\n" + "\n".join(files)
        except Exception as e: qdrant_status = f"Error checking Qdrant: {e}"
        return supabase_status, qdrant_status

    def handle_delete_file(filename):
        if not all((supabase, qdrant)): return "⚠️ Please set API keys first."
        if not filename: return "Please enter a filename to delete."
        filename = filename.strip()
        try:
            files_in_bucket = [f['name'] for f in supabase.storage.from_(BUCKET_NAME).list()]
            if filename not in files_in_bucket and sanitise_key(filename) not in files_in_bucket:
                return f"File '{filename}' not found in storage."
            filename_to_delete = filename if filename in files_in_bucket else sanitise_key(filename)
            return delete_file(supabase, qdrant, filename_to_delete)
        except Exception as e: return f"Error during deletion process: {e}"

    with gr.Blocks(theme=gr.themes.Soft(), title="RAG File Vault") as demo:
        gr.Markdown("# 📁 Scalable RAG App")
        gr.Markdown("An all-in-one app to upload, query, check storage and delete documents in a RAG pipeline.")

        with gr.Tab("0. Set API Keys"):
            google_api_key_input = gr.Textbox(label="Google API Key", type="password")
            supabase_url_input = gr.Textbox(label="Supabase URL")
            supabase_service_key_input = gr.Textbox(label="Supabase Service Key", type="password")
            qdrant_url_input = gr.Textbox(label="Qdrant URL")
            qdrant_api_key_input = gr.Textbox(label="Qdrant API Key", type="password")
            credentials_submit_btn = gr.Button("Initialize Clients", variant="primary")
            credentials_output = gr.Textbox(label="Status", interactive=False)

        with gr.Tab("1. Upload & Process"):
            file_uploader = gr.File(
                label="Upload Files",
                file_count="multiple",
                file_types=[".pdf", ".docx", ".txt", ".pptx", ".csv", ".xlsx", ".md", ".html"]
            )
            upload_button = gr.Button("Process Files", variant="primary")
            clear_upload_button = gr.Button("Clear")
            with gr.Row():
                upload_log = gr.Textbox(label="Processing Log", lines=10, interactive=False)
                extracted_preview = gr.Textbox(label="Extracted Text Preview", lines=10, interactive=False)

        with gr.Tab("2. Query Documents"):
            gr.Markdown("### Ask a Question")
            query_input = gr.Textbox(label="Your Question", placeholder="e.g., What are the laws of thermodynamics?")
            with gr.Row():
                query_button = gr.Button("Get Answer", variant="primary")
                clear_query_button = gr.Button("Clear")
            query_output = gr.Textbox(label="Model Answer", lines=8, interactive=False)
            retrieved_context = gr.Textbox(label="Retrieved Context", lines=10, interactive=False)

        with gr.Tab("3. Manage Storage"):
            gr.Markdown("### Check Cloud Storage Usage")
            check_storage_button = gr.Button("Check Usage", variant="primary")
            clear_storage_button = gr.Button("Clear")
            with gr.Row():
                supabase_status_output = gr.Textbox(label="Supabase Status", lines=15, interactive=False)
                qdrant_status_output = gr.Textbox(label="Qdrant Status", lines=15, interactive=False)

        with gr.Tab("4. Delete Files"):
            gr.Markdown("### Delete a File and its Embeddings")
            delete_filename_input = gr.Textbox(label="Filename to Delete", placeholder="Enter the exact filename, e.g., my-document.pdf")
            delete_button = gr.Button("Delete File Permanently", variant="stop")
            clear_delete_button = gr.Button("Clear")
            delete_log_output = gr.Textbox(label="Deletion Log", lines=5, interactive=False)

        # --- Event Handlers ---
        credentials_submit_btn.click(handle_credentials_submit, [google_api_key_input, supabase_url_input, supabase_service_key_input, qdrant_url_input, qdrant_api_key_input], [credentials_output, google_api_key_input, supabase_url_input, supabase_service_key_input, qdrant_url_input, qdrant_api_key_input])
        upload_button.click(handle_file_upload, [file_uploader], [upload_log, extracted_preview])
        clear_upload_button.click(lambda: (None, "", ""), None, [file_uploader, upload_log, extracted_preview])
        query_button.click(handle_rag_query, [query_input], [query_output, retrieved_context])
        clear_query_button.click(lambda: ("", "", ""), None, [query_input, query_output, retrieved_context])
        check_storage_button.click(handle_check_storage, None, [supabase_status_output, qdrant_status_output])
        clear_storage_button.click(lambda: ("", ""), None, [supabase_status_output, qdrant_status_output])
        delete_button.click(handle_delete_file, [delete_filename_input], [delete_log_output])
        clear_delete_button.click(lambda: ("", ""), None, [delete_filename_input, delete_log_output])
    return demo

# --- 8. Launch the App ---
if __name__ == "__main__":
    print("🚀 Launching RAG File Vault Gradio App...")
    app = create_gradio_app()
    app.launch(debug=True)
